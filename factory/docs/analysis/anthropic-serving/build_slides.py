from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
BASE="/Users/socrateshlapolosa/Development/health-service-idp/factory/docs/analysis/anthropic-serving"
DECK="/Users/socrateshlapolosa/Downloads/roadmap.pptx"
INK=C(0x1F,0x29,0x37); MUT=C(0x6B,0x72,0x80); BLUE=C(0x25,0x63,0xEB); TEAL=C(0x0D,0x94,0x88)
RED=C(0xC0,0x39,0x2B); GRN=C(0x15,0x80,0x3D); PUR=C(0x6D,0x28,0xD9); ORA=C(0xE0,0x8A,0x1E)
LGREY=C(0xF3,0xF4,0xF6); WHITE=C(0xFF,0xFF,0xFF); CARDLINE=C(0xD1,0xD5,0xDB)
prs=Presentation(DECK); BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def bar(s,color): 
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,I(0),I(0),I(0.16),I(7.5)); r.fill.solid(); r.fill.fore_color.rgb=color; r.line.fill.background()
def tb(s,l,t,w,h,runs,size=14,bold=False,color=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=4):
    b=s.shapes.add_textbox(I(l),I(t),I(w),I(h)); tf=b.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    if isinstance(runs,str): runs=[(runs,size,bold,color)]
    for i,item in enumerate(runs):
        txt,sz,bd,cl=item if len(item)==4 else (item[0],size,bold,color)
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.space_after=Pt(sp)
        r=p.add_run(); r.text=txt; f=r.font; f.size=Pt(sz); f.bold=bd; f.color.rgb=cl
    return b
def card(s,l,t,w,h,fill,line=CARDLINE,rad=True):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE,I(l),I(t),I(w),I(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(1); sh.shadow.inherit=False
    return sh
def bullets(s,l,t,w,h,items,size=12,color=INK):
    b=s.shapes.add_textbox(I(l),I(t),I(w),I(h)); tf=b.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(5)
        r=p.add_run(); r.text="•  "+it; r.font.size=Pt(size); r.font.color.rgb=color
    return b
def footer(s,txt="HSO · Agentic AI platform"):
    tb(s,0.16,7.12,8,0.35,txt,10,False,MUT)

# ---------- S1: section header ----------
s=slide(); bar(s,PUR)
tb(s,0.6,2.1,12,0.5,"AGENTIC AI",16,True,PUR)
tb(s,0.6,2.55,12,1.0,"Towards an Agentic AI Platform",40,True,INK)
tb(s,0.6,3.65,11.8,1.4,[
 ("We are adopting agentic AI. To scale it safely we need a platform — to HOST agentic solutions and to GOVERN them.",16,False,INK),
 ("A recent audit finding on Claude Code makes the governance need concrete. First: what is an agentic platform?",16,False,MUT)])
for i,(t,c) in enumerate([("① Host agentic solutions",TEAL),("② Govern agentic solutions",BLUE),("⚠ Audit finding — Claude Code",RED)]):
    card(s,0.6+i*4.05,5.4,3.8,0.7,LGREY,c); tb(s,0.75+i*4.05,5.55,3.6,0.45,t,13,True,c)
footer(s)

# ---------- S2: maturity ----------
s=slide(); bar(s,BLUE)
tb(s,0.5,0.35,12.3,0.8,"From AI to Agentic AI — the maturity",30,True,INK)
tb(s,0.5,1.05,12.3,0.5,"Agentic AI = autonomous agents that plan, use tools, remember and act to complete multi-step processes.",15,False,MUT)
s.shapes.add_picture(os.path.join(BASE,"maturity.png"),I(0.6),I(2.0),width=I(12.1))
tb(s,0.6,6.4,12.1,0.6,"Each stage builds on the last. An agentic platform is what makes the rightmost stage — Agentic AI — operable and safe at scale.",14,False,INK)
footer(s)

# ---------- S3: what is an agentic platform ----------
s=slide(); bar(s,TEAL)
tb(s,0.5,0.35,12.3,0.8,"What is an Agentic Platform?",30,True,INK)
tb(s,0.5,1.05,12.3,0.5,"The substrate to BUILD · RUN · MANAGE · GOVERN autonomous agents (per the Agentic AI framework).",15,False,MUT)
cols=[("Key Technologies",ORA,["LLMs","Transformers","Attention","Transfer learning","RAG / grounding"]),
      ("Agent Capabilities",BLUE,["Planning (ReAct/CoT)","Tool use & function calling","Memory (short/long)","Multi-agent collaboration","Self-reflection & recovery"]),
      ("Agent Management",TEAL,["Task scheduling","Rollback","Self-improvement","Feedback loops","Cost & resource mgmt"]),
      ("Outputs & Interfaces",PUR,["APIs","UIs / chat","Code & content gen","Speech (TTS / ASR)"]),
      ("Governance & Future",RED,["Identity (Entra)","Safety & guardrails","Memory governance","Observability & tracing","Risk mgmt · human-in-loop"])]
cw,gap,x0,y,h=2.42,0.12,0.5,1.7,4.6
for i,(t,c,items) in enumerate(cols):
    x=x0+i*(cw+gap); card(s,x,y,cw,h,WHITE,CARDLINE)
    hd=card(s,x,y,cw,0.55,c); tb(s,x+0.12,y+0.07,cw-0.2,0.45,t,12.5,True,WHITE)
    bullets(s,x+0.14,y+0.7,cw-0.24,h-0.8,items,11)
card(s,0.5,6.55,12.34,0.62,LGREY,TEAL)
tb(s,0.5,6.66,12.34,0.45,"An agentic platform provides the substrate to BUILD · RUN · MANAGE · GOVERN autonomous agents — reusing our existing landing zones.",14,True,INK,PP_ALIGN.CENTER)
footer(s)

# ---------- S4: why now — two imperatives ----------
s=slide(); bar(s,BLUE)
tb(s,0.5,0.35,12.3,0.8,"Why we need it now — two imperatives",30,True,INK)
# host card
card(s,0.5,1.55,6.05,3.6,WHITE,TEAL); card(s,0.5,1.55,6.05,0.62,TEAL); tb(s,0.7,1.68,5.8,0.45,"① HOST agentic solutions",16,True,WHITE)
tb(s,0.7,2.32,5.7,0.5,"A place to build & run agents, tools, orchestration and models — reuse our platform.",13,False,INK)
bullets(s,0.72,2.95,5.7,2.1,["Agent runtime — Foundry Agent Service","Tools via MCP","Orchestration & microservices (Container Apps / AKS)","Model serving (GPT-5.4 today · Claude next)","Scales our existing Azure landing zones"],12)
# govern card
card(s,6.8,1.55,6.05,3.6,WHITE,BLUE); card(s,6.8,1.55,6.05,0.62,BLUE); tb(s,7.0,1.68,5.8,0.45,"② GOVERN agentic solutions",16,True,WHITE)
tb(s,7.0,2.32,5.7,0.5,"Identity, policy, safety, observability and audit around every agent.",13,False,INK)
bullets(s,7.02,2.95,5.7,2.1,["Identity — Entra (incl. Agent 365 / Agent ID)","APIM AI Gateway — quota · cache · content safety","Guardrails — Azure Policy","Data & threat — Purview · Defender","Observability — Monitor / Log Analytics · human-in-loop"],12)
# audit callout
card(s,0.5,5.45,12.34,1.55,C(0xFD,0xEC,0xEA),RED)
tb(s,0.72,5.6,12,0.4,"⚠  Audit finding — Claude Code",15,True,RED)
tb(s,0.72,6.05,12,0.9,"Adopted direct-to-Anthropic SaaS: no enterprise identity (Entra), no API gateway/policy, no DLP or audit, and prompt + output data leaves the tenant.  →  governed hosting required.",13,False,INK)
footer(s)

# ---------- S5: as-is vs to-be ----------
s=slide(); bar(s,GRN)
tb(s,0.5,0.35,12.3,0.7,"Claude Code — As-Is vs To-Be",30,True,INK)
tb(s,0.5,1.0,12.3,0.45,"Same developer experience — moved behind enterprise identity, gateway and governance to close the audit finding.",14,False,MUT)
# left As-Is
card(s,0.5,1.55,6.05,0.5,RED); tb(s,0.7,1.62,5.8,0.4,"AS-IS  ·  direct to Anthropic   ❌ audit gap",14,True,WHITE)
s.shapes.add_picture(os.path.join(BASE,"view1-claude-code-direct.png"),I(0.55),I(2.15),width=I(5.95))
bullets(s,0.6,5.7,5.95,1.3,["Direct to Anthropic SaaS · API key (not Entra)","No APIM / no policy · no DLP / no audit","Prompt + output data egress outside the tenant"],11.5,RED)
# right To-Be
card(s,6.8,1.55,6.05,0.5,GRN); tb(s,7.0,1.62,5.8,0.4,"TO-BE  ·  via Foundry + APIM   ✅ finding closed",14,True,WHITE)
s.shapes.add_picture(os.path.join(BASE,"view2-logical-clean.png"),I(6.85),I(2.15),width=I(5.95))
bullets(s,6.9,5.7,5.95,1.3,["Entra-authenticated (corporate ID) → APIM → Foundry (Claude)","Content safety · quota · semantic cache","Defender · Purview · Policy · Monitor — private network, no public egress"],11.5,GRN)
footer(s)

prs.save(DECK)
print("saved:", DECK, "· total slides:", len(prs.slides.__iter__.__self__._sldIdLst))
